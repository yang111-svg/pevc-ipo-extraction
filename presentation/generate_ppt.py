# -*- coding: utf-8 -*-
"""Generate 黄山谷捷 presentation PPT — clean, minimalist, light-colored."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette (light, warm academic) ──
BG = RGBColor(0xFA, 0xFA, 0xF8)          # warm white
DARK = RGBColor(0x2D, 0x2D, 0x2D)         # near black
ACCENT = RGBColor(0x3B, 0x6F, 0x7D)       # teal accent
ACCENT_LIGHT = RGBColor(0xE8, 0xF0, 0xF2) # light teal
GRAY = RGBColor(0x78, 0x78, 0x78)         # gray text
RED = RGBColor(0xC0, 0x39, 0x2B)          # alert red
GREEN = RGBColor(0x27, 0xAE, 0x60)        # positive green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_title_bar(slide, title, subtitle=None):
    """Top bar with teal accent line + title"""
    # Accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(0.6), Inches(0.06), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    add_text_box(slide, 1.05, 0.55, 10, 0.6, title, font_size=32, bold=True, color=DARK)
    if subtitle:
        add_text_box(slide, 1.05, 1.05, 10, 0.4, subtitle, font_size=14, color=GRAY)
    # Divider line
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.01)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape2.line.fill.background()


def add_card(slide, left, top, width, height, title, content, title_size=14, content_size=11):
    """A rounded card with title and bullet content."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    shape.line.width = Pt(0.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = 'Microsoft YaHei'

    for line in content.strip().split('\n'):
        p = tf.add_paragraph()
        p.text = line.strip()
        p.font.size = Pt(content_size)
        p.font.color.rgb = DARK
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(2)


def add_footer(slide, text="杨苗鑫 · 黄山谷捷深度分析 · 2026.07"):
    add_text_box(slide, 0.8, 7.0, 11, 0.3, text, font_size=9, color=GRAY)


def add_table(slide, left, top, col_widths, headers, rows, font_size=10):
    """Add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_width = sum(col_widths)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top), Inches(table_width), Inches(0.35 * n_rows)
    )
    table = table_shape.table

    for j, w in enumerate(col_widths):
        table.columns[j].width = Inches(w)

    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════
# Slide 1: Title
# ═══════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, BG)

# Large accent block
shape = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(2.8), Inches(13.333), Inches(2.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide1, 1.5, 2.0, 10, 0.6, 'PEVC 招股书关键页抽取', font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 1.5, 2.9, 10, 0.8, '黄山谷捷 (301581) 深度分析', font_size=38, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 1.5, 3.7, 10, 0.6, '股权转让 · 吸收合并 · 出资额与比例并存 · 单位口径辨析', font_size=16, color=RGBColor(0xD0, 0xE8, 0xEC), alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 1.5, 5.5, 10, 0.4, '杨苗鑫    |    2026年7月', font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 1.5, 5.9, 10, 0.4, 'github.com/yang111-svg/pevc-ipo-extraction', font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 2: Progress Summary
# ═══════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, BG)
add_title_bar(slide2, '前三周进展总结', 'PEVC IPO 招股说明书股本变化事件自动抽取项目')

add_card(slide2, 0.8, 1.8, 3.6, 2.0, 'Week 1 · PDF获取与解析',
         '- 下载8家IPO公司PDF\n- 建立7步pipeline骨架\n- 项目结构初始化\n- 下载方法文档化')
add_card(slide2, 4.8, 1.8, 3.6, 2.0, 'Week 2 · 三表体系设计',
         '- 出资流水 + 股权快照 + 交叉核验\n- 16份Excel结构化输出\n- Schema定义(除转让外)\n- LLM Prompt v1')
add_card(slide2, 8.8, 1.8, 3.6, 2.0, 'Week 3 · Gold标注 + 自动化',
         '- 169条Gold记录\n- 事件分类过滤器(-31%噪声)\n- annotations_pdf 31页批注\n- 10步自动化管线')

# Core numbers
add_card(slide2, 0.8, 4.2, 3.6, 1.5, '核心技术指标',
         '• Gold记录: 169条\n• 覆盖: 8家公司\n• 管线步骤: 10步\n• 候选事件: 269个(黄山谷捷)')
add_card(slide2, 4.8, 4.2, 3.6, 1.5, '主要困难',
         '• PDF跨页表格合并\n• 候选块噪声(LLM前)\n• 股权转让vs增资边界\n• 单位口径混淆')
add_card(slide2, 8.8, 4.2, 3.6, 1.5, 'Week 4 · 方法升级',
         '• 代码定位章节(非LLM全文)\n• Schema: +TransferFlow\n• 规则候选过滤\n• MinerU Markdown')

add_footer(slide2)

# ═══════════════════════════════════════════
# Slide 3: Company Timeline
# ═══════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, BG)
add_title_bar(slide3, '黄山谷捷 · 股本变化时间线', '301581 · 黄山 · 散热基板 · 从设立到股份公司 (2012-2022)')

timeline_data = [
    ['2012/06', '设立', '1,000万', '—', '—', '昆山谷捷100%出资'],
    ['2021/04/03', '第一次股权转让', '1,000万', '—', '0元(零对价)', '昆山谷捷→黄山供销78%/张11%/周11%'],
    ['2021/04/07', '吸收合并', '1,000→1,200万', '+200万', '0元', '吸收合并昆山谷捷，等比例增资'],
    ['2021/11/11', '第一次增资', '1,200→1,714.2857万', '+514.2857万', '22.42元/注资', '赛格高技术+上汽科技(PE/产业资本)'],
    ['2022/05/26', '第二次增资', '1,714.2857→1,804.5113万', '+90.2256万', '23.55元/注资', '黄山佳捷(员工持股平台)'],
    ['2022/09/13', '整体变更', '→6,000万股', '—', '1.0元/股(面值)', '净资产20,152.51万折股，各股东持股比不变'],
]

add_table(slide3, 0.8, 1.8,
          [1.5, 1.5, 2.2, 1.3, 1.7, 3.5],
          ['时间', '事件类型', '注册资本变化', '新增', '价格/对价', '关键信息'],
          timeline_data, font_size=10)

add_card(slide3, 0.8, 4.2, 11.7, 2.3, '关键事件特征',
         '• 股权转让(零对价) → 同一控制下架构调整，非市场化交易，须人工判断 event_type\n'
         '• 吸收合并(非现金增资) → 以被合并方净资产出资，subscription_amount=0 但 subscription_qty=200\n'
         '• 第一次增资(PE定价) → 赛格高技术 462.8571万元 × 22.42元 = 约10,377万元现金投入\n'
         '• 第二次增资(员工激励) → 员工持股平台，价格略高于PE(23.55 vs 22.42)，合理\n'
         '• 整体变更 → 净资产按1:0.29773折股，6,000万股 = 面值1元/股\n'
         '• 股东数量变化: 1→3→3→5→6→6 (在同一年内反复变化)')

add_footer(slide3)

# ═══════════════════════════════════════════
# Slide 4: PDF Evidence (page 1)
# ═══════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, BG)
add_title_bar(slide4, 'PDF 原文证据 (1/2)', '黄山谷捷招股说明书 第41-42页 · 关键数字来源')

# Evidence 1
add_card(slide4, 0.8, 1.8, 5.8, 2.3, '证据1: 报告期内股本变化总览 (PDF第41页, 1-1-40)',
         '招股书用时间线图清晰展示6个关键事件:\n'
         '  设立(2012/06) → 股权转让(2021/04) → 吸收合并(2021/04)\n'
         '  → 第一次增资(2021/11) → 第二次增资(2022/05) → 整体变更(2022/09)\n\n'
         '🟢 PDF直接披露: 事件顺序、各阶段注册资本\n'
         '🟡 须判断: "零对价转让"的理由 → PDF披露为"同一股权结构下的公司架构调整"')

# Evidence 2
add_card(slide4, 7.0, 1.8, 5.5, 2.3, '证据2: 第一次股权转让后股权结构 (PDF第42页, 1-1-41)',
         '股权转让后结构表:\n'
         '  黄山供销集团   780.00万元    78.00%\n'
         '  张俊武          110.00万元    11.00%\n'
         '  周斌            110.00万元    11.00%\n'
         '  合计           1,000.00万元   100.00%\n\n'
         '🔴 单位口径关键点:\n'
         '  "78" 在表中既可以是百分比(78.00%), 也可以是出资额(780万)')

# Evidence 3
add_card(slide4, 0.8, 4.4, 5.8, 2.3, '证据3: 吸收合并前后对照表 (PDF第43-44页, 1-1-42/43)',
         '三列表: 谷捷有限(合并前) | 昆山谷捷 | 谷捷有限(合并后)\n'
         '  黄山供销: 780 → (+156) → 936 (保持78%)\n'
         '  张俊武:   110 → (+22)  → 132 (保持11%)\n'
         '  周斌:     110 → (+22)  → 132 (保持11%)\n'
         '  合计:    1,000 → (+200) → 1,200\n\n'
         '🟢 PDF直接披露: 吸收合并前后出资额和比例\n'
         '🟡 须判断: 增资的200万 = "被合并方账面净资产" → 非现金出资')

add_footer(slide4)

# ═══════════════════════════════════════════
# Slide 5: PDF Evidence (page 2) + Gold Table
# ═══════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, BG)
add_title_bar(slide5, 'PDF 原文证据 (2/2) + Gold 表示例', '第44-45页 · 第一次增资定价与对赌协议')

add_card(slide5, 0.8, 1.8, 5.8, 2.3, '证据4: 第一次增资完成后的股权结构 (PDF第45页, 1-1-44)',
         '增资后结构:\n'
         '  黄山供销集团   936.0000万元    54.60%\n'
         '  赛格高技术      462.8571万元    27.00%   ← PE新进入\n'
         '  张俊武          132.0000万元     7.70%\n'
         '  周斌            132.0000万元     7.70%\n'
         '  上汽科技         51.4286万元     3.00%   ← 产业资本\n'
         '  合计           1,714.2857万元   100.00%\n\n'
         '🔴 直接披露 vs 须计算: 增资价格为22.42元/注册资本(直接披露)\n'
         '   赛格认购金额 = 462.8571 × 22.42 ≈ 10,377万元(须计算)')

add_card(slide5, 7.0, 1.8, 5.5, 2.3, 'Gold 表示例: 将PDF信息映射为结构化字段',
         '① subscription_flow (增资事件):\n'
         '  investor_name: "深圳赛格高技术投资股份有限公司"\n'
         '  investor_type: "PE"\n'
         '  subscription_qty_wan: 462.8571 ← PDF直接\n'
         '  subscription_price: 22.42      ← PDF直接\n'
         '  subscription_amount_wan: 10377  ← 须计算\n\n'
         '② transfer_flow (股权转让):\n'
         '  transferor: "昆山谷捷" → transferee: "黄山供销集团"\n'
         '  transfer_type: "零对价转让/同一控制下架构调整" ← 须判断\n'
         '  unit_issue_flag: true ← 78% vs 780万元')

# Field classification
add_card(slide5, 0.8, 4.4, 11.7, 2.3, '字段分类: 直接披露 vs 须计算 vs 须判断',
         '🟢 直接披露: 股东名称、持股比例(%)、出资额(万元)、增资价格(元/注册资本)、事件日期\n'
         '🟡 须计算: 认购总金额(=认购股数×价格)、每股单价(=总价÷股数)、总股本加总验证\n'
         '🔴 须判断: 事件类型(增资/转让/吸收合并)、零对价转让的定价合理性、单位口径(78%/780万元/78万元出资额)')
add_footer(slide5)

# ═══════════════════════════════════════════
# Slide 6: Automation Pipeline + Failure Points
# ═══════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, BG)
add_title_bar(slide6, '自动化流程与失败点分析', 'Week 4 改进后的 pipeline · 代码定位 + 规则过滤 + LLM 提取')

add_card(slide6, 0.8, 1.8, 7.5, 2.0, '改进后的自动化流程 (核心变化: 代码做定位，不是LLM)',
         'PDF → [PyMuPDF] 逐页文本 → [代码] 正则定位29个章节块\n'
         '  → [代码] 规则提取269个候选事件(金额/日期/股东/比例)\n'
         '  → [代码] 质量分级(high=76, medium=6, low=187)\n'
         '  → [代码] 专项检测(吸收合并×17, 单位混淆)\n'
         '  → [LLM] 仅处理高质量候选块 → 结构化JSON\n'
         '  → [代码] Schema校验 + Cross-Check')

add_card(slide6, 8.7, 1.8, 3.8, 2.0, 'Week 3 vs Week 4',
         'Week 3 (旧):\n'
         '• 全文→LLM (噪声重)\n'
         '• Schema缺TransferFlow\n'
         '• 无_raw原始值\n'
         '• 无unit_issue标记\n\n'
         'Week 4 (新):\n'
         '• 代码定位+规则过滤\n'
         '• TransferFlow独立类型\n'
         '• 原文值保留\n'
         '• 单位歧义自动标记')

fail_data = [
    ['PDF解析', '跨页表格合并', '中', '规则+人工补全'],
    ['章节定位', '标题变体多(29块中部分冗余)', '低', '增加关键词覆盖+层级匹配'],
    ['事件分类', '增资vs转让vs合并边界', '高', '新增事件类型判定规则'],
    ['单位识别', '78/780 单位混淆', '高', 'unit_issue_flag + 人工复核'],
    ['LLM提取', '非事件文本混入结果', '中', '前置规则过滤(↓31%噪声)'],
    ['交叉校验', '流量≠存量(舍入误差)', '低', '允许±0.01%容差'],
]
add_table(slide6, 0.8, 4.2,
          [2.0, 3.0, 0.8, 3.0],
          ['失败环节', '具体表现', '影响', '应对方案'],
          fail_data, font_size=10)

add_footer(slide6)

# ═══════════════════════════════════════════
# Slide 7: Special Difficulty + Commonality
# ═══════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, BG)
add_title_bar(slide7, '黄山谷捷特殊难点 & 跨公司共性', '从个案提炼方法 · IPO招股书股本变化披露的通用特征')

add_card(slide7, 0.8, 1.8, 5.8, 3.5, '黄山谷捷 5大特殊难点',
         '① 单位口径并存: PDF同时出现"78%"(比例)和"780万元"(出资额)\n'
         '   → 自动化抽取易混淆: 78→取比例还是取出资额?\n\n'
         '② 零对价转让: 昆山谷捷→黄山供销/张俊武/周斌，对价0元\n'
         '   → 须判断是"架构调整"而非市场交易\n\n'
         '③ 非现金增资: 吸收合并增资200万，以账面净资产出资\n'
         '   → subscription_amount=0 但 subscription_qty=200\n\n'
         '④ 同一控制下重组: 股权转让+吸收合并+增资密集发生\n'
         '   → 事件分组和时序容易错排\n\n'
         '⑤ 股东数量反复: 1→3→3→5→6→6 (同一年内)\n'
         '   → 需要正确追溯每步股东构成变化')

add_card(slide7, 7.0, 1.8, 5.5, 3.5, 'IPO 招股书股本变化披露的共性',
         '共性1 · 固定章节结构:\n'
         '  "发行人基本情况" → "股本和股东变化情况"\n'
         '  → 各事件按时间线逐一披露\n\n'
         '共性2 · 关键信息来源:\n'
         '  • 历史沿革章节的"出资额+比例"双列表\n'
         '  • 发行前股东结构表(最终时点快照)\n'
         '  • 对赌协议/股东特殊权利条款\n\n'
         '共性3 · 通用抽取难点:\n'
         '  • 金额单位(万元/元/万股)需统一转换\n'
         '  • 跨页表格需合并\n'
         '  • 增资与转让相邻易混淆\n'
         '  • 同一主体多次投资需去重\n\n'
         '共性4 · 自动化策略:\n'
         '  先找章节→再找表格→提取事件→填字段')

add_card(slide7, 0.8, 5.6, 11.7, 1.1, '如果换一家公司(如云汉芯城/赛分科技)，应该先去招股书第4-5节找"历史沿革"，从中提取出资额+股权结构表，然后按时间线逐一核对增资和转让事件',
         '', title_size=12, content_size=12)
add_footer(slide7)

# ═══════════════════════════════════════════
# Slide 8: Future Plan
# ═══════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, BG)
add_title_bar(slide8, '下一步改进计划', 'Gold 深化 · 自动化优化 · 方法推广')

add_card(slide8, 0.8, 1.8, 3.6, 2.8, 'Gold 标准深化 🔴 高优',
         '1. 补齐8家公司完整gold链\n'
         '   (设立→每次变更→当前快照)\n'
         '2. 每条gold加unit字段显式标注\n'
         '   万元/万股/股/元/美元\n'
         '3. review_queue闭环\n'
         '   人工确认PDF异常数据\n'
         '4. 增量gold与存量gold互相核验\n'
         '   流量(事件) = 存量变化(快照差)\n'
         '5. 目标: 每公司t0快照+全部事件+最终快照')

add_card(slide8, 4.8, 1.8, 3.6, 2.8, '自动化流程优化 🔴 高优',
         '1. 单位识别器\n'
         '   检测"出资额(万元)"vs"持股数量(股)"表头\n'
         '   自动标注单位 → 精确转换\n'
         '2. 事件分类器阈值调优\n'
         '   用gold数据反馈自动选最优阈值\n'
         '3. 零对价/非现金交易检测\n'
         '   关键词: "零对价""以净资产出资"\n'
         '4. 跨页表格合并\n'
         '   检测表格列结构一致性\n'
         '5. MinerU Markdown模式\n'
         '   保留表格结构 → 更精准定位')

add_card(slide8, 8.8, 1.8, 3.6, 2.8, '方法与推广 🟡 中优',
         '1. 1→8 家标准化\n'
         '   黄山谷捷试点 → 推广其余7家\n'
         '2. Prompt sensitivity analysis\n'
         '   用gold数据自动选最优prompt\n'
         '3. 定量评估体系\n'
         '   recall/precision/F1 per 公司\n'
         '   auto vs gold row_match 对比\n'
         '4. 最小可复现pipeline\n'
         '   python run_all.py 一键全流程\n'
         '5. 汇报与分享\n'
         '   把方法写清楚，让其他同学能复现')

add_card(slide8, 0.8, 4.9, 11.7, 1.8, '核心目标',
         '🎯 从"机械完成任务" → "理解问题结构"\n'
         '   让听的人明白: 如果换一家类似公司，应该先去招股书哪里找信息、哪些数字可以直接抽、哪些数字必须人工判断、自动化流程在哪里加规则/加LLM/加人工复核',
         title_size=13, content_size=13)
add_footer(slide8)


# ── Save ──
output_path = 'C:/Users/HP/Desktop/黄山谷捷_深度分析_杨苗鑫.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
